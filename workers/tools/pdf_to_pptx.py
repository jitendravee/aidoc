from __future__ import annotations
import os
import fitz
from pptx import Presentation
from pptx.util import Inches
from workers.tools.base import ToolInput, ToolOutput


class PdfToPptxInput(ToolInput):
    input_path: str
    output_path: str
    dpi: int = 150


class PdfToPptxOutput(ToolOutput):
    output_path: str
    diff_summary: str
    page_count: int


def run(params: PdfToPptxInput) -> PdfToPptxOutput:
    doc = fitz.open(params.input_path)
    zoom = params.dpi / 72
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    tmp_images = []

    for page in doc:
        pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
        img_path = f"{params.output_path}.tmp_{page.number}.jpg"
        pix.save(img_path)
        tmp_images.append(img_path)
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height) # type: ignore

    page_count = len(doc)
    doc.close()
    prs.save(params.output_path)
    for p in tmp_images:
        os.remove(p)

    return PdfToPptxOutput(
        output_path=params.output_path,
        diff_summary=f"Converted {page_count} page(s) to a PowerPoint presentation.",
        page_count=page_count,
    )